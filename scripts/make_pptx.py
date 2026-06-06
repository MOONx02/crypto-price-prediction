#!/usr/bin/env python3
"""
Generate a PowerPoint (.pptx) from a list of slides.
Usage:
  python scripts/make_pptx.py                    # uses DEFAULT_SLIDES below
  python scripts/make_pptx.py -o my_deck.pptx   # custom output path

Or import and call make_pptx(slides, path) with:
  slides = [ ("Title", ["bullet1", "bullet2", ...]), ... ]
"""
from pathlib import Path
import argparse

try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:
    print("Install python-pptx: pip install python-pptx")
    raise SystemExit(1)

DEFAULT_SLIDES = [
    ("Title", ["Subtitle or date"]),
    ("Slide 2", ["Point one", "Point two", "Point three"]),
    ("Conclusion", ["Summary point"]),
]


def make_pptx(slides, out_path: str | Path) -> Path:
    """Create a .pptx with title slide + content slides. Returns path."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title, *rest = slides[0] if slides else ("Title",)
    slide.shapes.title.text = title
    if rest and rest[0]:
        slide.placeholders[1].text = rest[0]

    # Content slides
    layout = prs.slide_layouts[1]
    for item in slides[1:]:
        title = item[0] if isinstance(item[0], str) else item[0]
        bullets = item[1] if len(item) > 1 else []
        if isinstance(bullets, str):
            bullets = [bullets]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame
        for b in bullets:
            p = tf.add_paragraph()
            p.text = str(b)
            p.level = 0

    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


def main():
    p = argparse.ArgumentParser(description="Generate .pptx from slide list")
    p.add_argument("-o", "--output", default="output.pptx", help="Output .pptx path")
    p.add_argument("--default", action="store_true", help="Use built-in default slides")
    args = p.parse_args()
    slides = DEFAULT_SLIDES
    path = make_pptx(slides, args.output)
    print(f"Saved: {path.resolve()}")


if __name__ == "__main__":
    main()
